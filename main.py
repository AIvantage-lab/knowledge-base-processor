from fastapi import FastAPI, HTTPException, BackgroundTasks
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from typing import Dict, Any, List, Optional
import PyPDF2
import pdfplumber
from docx import Document
import openpyxl
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image
import io
import tempfile
import os
from langdetect import detect
import httpx
from openai import OpenAI
from supabase import create_client, Client
import hashlib
import time
import json
import re
import asyncio

app = FastAPI(title="Knowledge Base Processor Service")

# Configuración de clientes
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
SUPABASE_URL = os.getenv("SUPABASE_URL", "https://hsoagaoxuaspkdptfgou.supabase.co")
SUPABASE_SERVICE_KEY = os.getenv("SUPABASE_SERVICE_KEY")


# ============================================================
# DISTRIBUCIÓN DE PROGRESO POR ETAPAS
# ============================================================
# 0-5%:    Descarga del archivo
# 5-25%:   Extracción de texto (incluye OCR si es necesario)
# 25-30%:  Chunking y preparación
# 30-95%:  Generación de embeddings e inserción
# 95-100%: Extracción de metadata con LLM y finalización
# ============================================================

PROGRESS_STAGES = {
    "download_start": 0.0,
    "download_complete": 0.05,
    "extraction_start": 0.05,
    "extraction_complete": 0.25,
    "chunking_start": 0.25,
    "chunking_complete": 0.30,
    "embeddings_start": 0.30,
    "embeddings_complete": 0.95,
    "metadata_start": 0.95,
    "metadata_complete": 1.0
}


class ProcessDocumentRequest(BaseModel):
    file_url: str
    library_id: str
    file_name: str
    subject: str
    subject_id: Optional[str] = None
    authors: Optional[List[str]] = None
    title: Optional[str] = None
    publication_date: Optional[str] = None


def get_supabase_client() -> Client:
    """Crea y retorna un cliente de Supabase"""
    return create_client(SUPABASE_URL, SUPABASE_SERVICE_KEY)


def update_progress(supabase: Client, library_id: str, percentage: float, status: str = "processing"):
    """
    Actualiza el progreso en AIvantage_library.
    
    Args:
        supabase: Cliente de Supabase
        library_id: ID del documento en AIvantage_library
        percentage: Porcentaje de progreso (0.0 a 1.0)
        status: Estado del procesamiento
    """
    try:
        supabase.table("AIvantage_library").update({
            "loading_percentage": round(percentage, 3),
            "status": status
        }).eq("id", library_id).execute()
        print(f"[Progreso] {percentage*100:.1f}% - {status}")
    except Exception as e:
        print(f"[Error actualizando progreso]: {str(e)}")


def calculate_embeddings_progress(current_chunk: int, total_chunks: int) -> float:
    """
    Calcula el progreso dentro de la fase de embeddings (30-95%).
    
    Args:
        current_chunk: Chunk actual procesado
        total_chunks: Total de chunks
    
    Returns:
        Porcentaje de progreso total (entre 0.30 y 0.95)
    """
    embeddings_range = PROGRESS_STAGES["embeddings_complete"] - PROGRESS_STAGES["embeddings_start"]
    chunk_progress = current_chunk / total_chunks
    return PROGRESS_STAGES["embeddings_start"] + (embeddings_range * chunk_progress)


@app.get("/")
async def health_check():
    return {"status": "healthy", "service": "knowledge-base-processor"}


@app.post("/process")
async def process_document(request: ProcessDocumentRequest, background_tasks: BackgroundTasks):
    """
    Endpoint asíncrono para procesar documentos.
    
    Responde inmediatamente con status "queued" y procesa en background.
    El progreso se actualiza en tiempo real en AIvantage_library.
    """
    try:
        # Validar configuración antes de encolar
        if not OPENAI_API_KEY:
            raise HTTPException(status_code=500, detail="OPENAI_API_KEY no configurada")
        if not SUPABASE_SERVICE_KEY:
            raise HTTPException(status_code=500, detail="SUPABASE_SERVICE_KEY no configurada")
        
        # Inicializar Supabase para actualizar status inicial
        supabase = get_supabase_client()
        
        # Marcar como "queued" inmediatamente
        supabase.table("AIvantage_library").update({
            "status": "queued",
            "loading_percentage": 0.0
        }).eq("id", request.library_id).execute()
        
        # Encolar el procesamiento en background
        background_tasks.add_task(
            process_document_background,
            request.file_url,
            request.library_id,
            request.file_name,
            request.subject,
            request.subject_id,
            request.authors,
            request.title,
            request.publication_date
        )
        
        # Responder inmediatamente
        return {
            "success": True,
            "status": "queued",
            "message": "Documento encolado para procesamiento",
            "library_id": request.library_id
        }
        
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error en process_document: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


async def process_document_background(
    file_url: str,
    library_id: str,
    file_name: str,
    subject: str,
    subject_id: Optional[str],
    authors: Optional[List[str]],
    title: Optional[str],
    publication_date: Optional[str]
):
    """
    Procesa un documento en background.
    
    Etapas:
    1. Descarga el documento (0-5%)
    2. Extrae el texto (5-25%)
    3. Crea chunks (25-30%)
    4. Genera embeddings e inserta (30-95%)
    5. Extrae metadata con LLM (95-100%)
    """
    supabase = None
    record_id = None
    
    try:
        # Inicializar clientes
        openai_client = OpenAI(api_key=OPENAI_API_KEY)
        supabase = get_supabase_client()
        
        # ============================================================
        # ETAPA 1: DESCARGA (0-5%)
        # ============================================================
        update_progress(supabase, library_id, PROGRESS_STAGES["download_start"], "downloading")
        
        print(f"[Etapa 1] Descargando documento: {file_url}")
        async with httpx.AsyncClient(timeout=300.0) as client:
            response = await client.get(file_url)
            if response.status_code != 200:
                update_progress(supabase, library_id, 0, "error")
                print(f"[Error] No se pudo descargar el archivo: {response.status_code}")
                return
            file_content = response.content
        
        update_progress(supabase, library_id, PROGRESS_STAGES["download_complete"], "processing")
        print(f"[Etapa 1] Documento descargado: {len(file_content)} bytes")
        
        # ============================================================
        # ETAPA 2: EXTRACCIÓN DE TEXTO (5-25%)
        # ============================================================
        update_progress(supabase, library_id, PROGRESS_STAGES["extraction_start"], "extracting")
        
        file_extension = file_name.split('.')[-1].lower()
        print(f"[Etapa 2] Extrayendo texto de archivo {file_extension}...")
        
        if file_extension == 'pdf':
            extraction_result = await extract_pdf_with_progress(file_content, supabase, library_id)
        elif file_extension in ['docx', 'doc']:
            extraction_result = await extract_docx(file_content)
        elif file_extension in ['xlsx', 'xls']:
            extraction_result = await extract_excel(file_content)
        elif file_extension in ['pptx', 'ppt']:
            extraction_result = await extract_pptx(file_content)
        elif file_extension in ['png', 'jpg', 'jpeg']:
            extraction_result = await extract_image(file_content)
        elif file_extension in ['txt', 'md']:
            extraction_result = await extract_text(file_content)
        else:
            update_progress(supabase, library_id, 0, "error")
            print(f"[Error] Formato no soportado: {file_extension}")
            return
        
        extracted_text = extraction_result.get('text', '')
        if not extracted_text or len(extracted_text.strip()) < 50:
            update_progress(supabase, library_id, 0, "error")
            print("[Error] No se pudo extraer texto suficiente del documento")
            return
        
        update_progress(supabase, library_id, PROGRESS_STAGES["extraction_complete"], "processing")
        print(f"[Etapa 2] Texto extraído: {len(extracted_text)} caracteres")
        
        # Detectar idioma
        try:
            language = detect(extracted_text[:1000])
        except:
            language = 'unknown'
        
        # ============================================================
        # ETAPA 3: CHUNKING (25-30%)
        # ============================================================
        update_progress(supabase, library_id, PROGRESS_STAGES["chunking_start"], "chunking")
        
        print("[Etapa 3] Creando chunks...")
        chunks = create_smart_chunks(extracted_text, chunk_size=4000, overlap=400)
        
        # Generar hash del documento
        doc_hash = hashlib.md5(file_content).hexdigest()
        
        # Insertar en record_manager_v2
        record_data = {
            "file_path": file_url,
            "file_name": file_name,
            "content_hash": doc_hash,
            "status": "processing",
            "total_chunks": len(chunks),
            "processed_chunks": 0,
            "library_id": library_id
        }
        
        record_result = supabase.table("record_manager_v2").insert(record_data).execute()
        
        if not record_result.data:
            update_progress(supabase, library_id, 0, "error")
            print("[Error] No se pudo crear registro en record_manager_v2")
            return
        
        record_id = record_result.data[0]['id']
        
        update_progress(supabase, library_id, PROGRESS_STAGES["chunking_complete"], "processing")
        print(f"[Etapa 3] Chunks creados: {len(chunks)}, Record ID: {record_id}")
        
        # ============================================================
        # ETAPA 4: EMBEDDINGS E INSERCIÓN (30-95%)
        # ============================================================
        update_progress(supabase, library_id, PROGRESS_STAGES["embeddings_start"], "embedding")
        
        print("[Etapa 4] Generando embeddings...")
        batch_size = 20
        total_vectors_inserted = 0
        
        for i in range(0, len(chunks), batch_size):
            batch_chunks = chunks[i:i + batch_size]
            texts_to_embed = [chunk['content'] for chunk in batch_chunks]
            
            try:
                embedding_response = openai_client.embeddings.create(
                    model="text-embedding-3-small",
                    input=texts_to_embed
                )
                
                # Preparar datos para insertar en documents_v2
                documents_to_insert = []
                
                for j, (chunk, embedding_data) in enumerate(zip(batch_chunks, embedding_response.data)):
                    chunk_index = i + j
                    
                    doc_record = {
                        "content": chunk['content'],
                        "metadata": {
                            "file_name": file_name,
                            "subject": subject,
                            "subject_id": subject_id,
                            "authors": authors,
                            "title": title or file_name,
                            "publication_date": publication_date,
                            "chunk_index": chunk_index,
                            "total_chunks": len(chunks),
                            "language": language,
                            "library_id": library_id,
                            "source": "knowledge_base"
                        },
                        "embedding": embedding_data.embedding,
                        "record_id": record_id
                    }
                    documents_to_insert.append(doc_record)
                
                # Insertar en documents_v2
                if documents_to_insert:
                    supabase.table("documents_v2").insert(documents_to_insert).execute()
                    total_vectors_inserted += len(documents_to_insert)
                
                # Actualizar progreso en record_manager_v2
                progress_chunks = min(i + batch_size, len(chunks))
                supabase.table("record_manager_v2").update({
                    "processed_chunks": progress_chunks
                }).eq("id", record_id).execute()
                
                # Actualizar progreso en AIvantage_library (30-95%)
                current_progress = calculate_embeddings_progress(progress_chunks, len(chunks))
                update_progress(supabase, library_id, current_progress, "embedding")
                
                print(f"[Etapa 4] Lote procesado: {progress_chunks}/{len(chunks)} chunks")
                
                # Pequeña pausa para evitar rate limits
                if i + batch_size < len(chunks):
                    await asyncio.sleep(0.5)
                    
            except Exception as e:
                print(f"[Error] Error procesando lote {i}: {str(e)}")
                continue
        
        # Actualizar estado en record_manager_v2
        supabase.table("record_manager_v2").update({
            "status": "completed",
            "processed_chunks": len(chunks)
        }).eq("id", record_id).execute()
        
        update_progress(supabase, library_id, PROGRESS_STAGES["embeddings_complete"], "finalizing")
        print(f"[Etapa 4] Embeddings completados. Vectores insertados: {total_vectors_inserted}")
        
        # ============================================================
        # ETAPA 5: EXTRACCIÓN DE METADATA CON LLM (95-100%)
        # ============================================================
        update_progress(supabase, library_id, PROGRESS_STAGES["metadata_start"], "extracting_metadata")
        
        print("[Etapa 5] Extrayendo metadata con LLM...")
        
        try:
            # Extraer metadata de los chunks
            extracted_metadata = await extract_metadata_with_llm(
                chunks=chunks,
                openai_client=openai_client,
                file_name=file_name
            )
            
            # Obtener datos actuales de AIvantage_library para comparar
            current_library_data = supabase.table("AIvantage_library").select(
                "title, authors, publication_date"
            ).eq("id", library_id).execute()
            
            if current_library_data.data:
                current_data = current_library_data.data[0]
                
                # Comparar y actualizar si corresponde
                metadata_result = await compare_and_update_metadata(
                    extracted_metadata=extracted_metadata,
                    library_id=library_id,
                    current_data=current_data,
                    supabase=supabase
                )
                
                print(f"[Etapa 5] Metadata actualizada: {metadata_result.get('changes_made', [])}")
                
        except Exception as e:
            print(f"[Etapa 5] Error en extracción de metadata (no crítico): {str(e)}")
        
        # ============================================================
        # FINALIZACIÓN
        # ============================================================
        update_progress(supabase, library_id, PROGRESS_STAGES["metadata_complete"], "processed")
        
        print(f"[Completado] Documento procesado exitosamente.")
        print(f"  - Library ID: {library_id}")
        print(f"  - Record ID: {record_id}")
        print(f"  - Total chunks: {len(chunks)}")
        print(f"  - Total vectores: {total_vectors_inserted}")
        
    except Exception as e:
        print(f"[Error Fatal] Error en process_document_background: {str(e)}")
        
        # Intentar actualizar status a error
        if supabase:
            try:
                supabase.table("AIvantage_library").update({
                    "status": "error",
                    "loading_percentage": 0
                }).eq("id", library_id).execute()
                
                if record_id:
                    supabase.table("record_manager_v2").update({
                        "status": "error"
                    }).eq("id", record_id).execute()
            except:
                pass


# ============================================================
# FUNCIÓN DE EXTRACCIÓN DE PDF CON PROGRESO
# ============================================================

async def extract_pdf_with_progress(content: bytes, supabase: Client, library_id: str) -> Dict[str, Any]:
    """
    Extrae texto de PDF con reporte de progreso.
    Incluye sub-progreso para OCR (5-25% del total).
    """
    result = {'text': '', 'pages': 0}
    
    # Sub-etapas dentro de extracción (5-25%)
    extraction_start = PROGRESS_STAGES["extraction_start"]
    extraction_end = PROGRESS_STAGES["extraction_complete"]
    extraction_range = extraction_end - extraction_start
    
    try:
        # Intentar extracción normal primero (rápido)
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        result['pages'] = len(pdf_reader.pages)
        
        text_parts = []
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"--- Página {page_num + 1} ---\n{page_text}")
        
        result['text'] = '\n\n'.join(text_parts)
        
        # Si hay texto suficiente, terminamos (actualizar a 25%)
        if result['text'].strip() and len(result['text']) >= 100:
            return result
        
        # Intentar con pdfplumber
        print("[Extracción] Intentando con pdfplumber...")
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
            tmp_file.write(content)
            tmp_path = tmp_file.name
        
        with pdfplumber.open(tmp_path) as pdf:
            result['pages'] = len(pdf.pages)
            text_parts = []
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(f"--- Página {i + 1} ---\n{page_text}")
            result['text'] = '\n\n'.join(text_parts)
        
        os.unlink(tmp_path)
        
        # Si hay texto suficiente, terminamos
        if result['text'].strip() and len(result['text']) >= 100:
            return result
        
        # Si no hay texto, usar OCR (proceso largo)
        print("[Extracción] Iniciando OCR (esto puede tomar varios minutos)...")
        result['text'] = await extract_pdf_with_ocr_progress(content, supabase, library_id)
            
    except Exception as e:
        print(f"[Error] Error extrayendo PDF: {e}")
        print("[Extracción] Intentando OCR como fallback...")
        result['text'] = await extract_pdf_with_ocr_progress(content, supabase, library_id)
    
    return result


async def extract_pdf_with_ocr_progress(content: bytes, supabase: Client, library_id: str) -> str:
    """
    Extrae texto de PDF usando OCR con reporte de progreso por página.
    El OCR ocupa la fase 5-25% del progreso total.
    """
    try:
        # Convertir PDF a imágenes
        images = convert_from_bytes(content, dpi=200)
        total_pages = len(images)
        
        print(f"[OCR] Procesando {total_pages} páginas...")
        
        # Calcular progreso dentro del rango de extracción (5-25%)
        extraction_start = PROGRESS_STAGES["extraction_start"]
        extraction_end = PROGRESS_STAGES["extraction_complete"]
        extraction_range = extraction_end - extraction_start
        
        text_parts = []
        
        for i, image in enumerate(images):
            # Procesar página con OCR
            text = pytesseract.image_to_string(image, lang='spa+eng')
            text_parts.append(f"--- Página {i + 1} (OCR) ---\n{text}")
            
            # Actualizar progreso cada página
            page_progress = (i + 1) / total_pages
            current_progress = extraction_start + (extraction_range * page_progress)
            update_progress(supabase, library_id, current_progress, "extracting_ocr")
            
            # Log cada 10 páginas o en la última
            if (i + 1) % 10 == 0 or (i + 1) == total_pages:
                print(f"[OCR] Página {i + 1}/{total_pages} procesada ({current_progress*100:.1f}%)")
            
            # Pequeña pausa para no bloquear el event loop
            await asyncio.sleep(0.01)
        
        return '\n\n'.join(text_parts)
        
    except Exception as e:
        print(f"[Error] Error en OCR: {str(e)}")
        return f"Error en OCR: {str(e)}"


# ============================================================
# FUNCIÓN DE FORMATEO DE FECHA
# ============================================================

def format_publication_date(date_string: str) -> str:
    """
    Convierte una fecha en varios formatos posibles a DD/MM/YYYY.
    Si solo hay año, retorna 01/01/YYYY.
    Si no puede parsear, retorna 'NO_IDENTIFICADO'.
    """
    if not date_string or date_string == "NO_ENCONTRADO":
        return "NO_IDENTIFICADO"
    
    # Limpiar el string
    date_string = date_string.strip()
    
    # Diccionario de meses en español e inglés
    months = {
        'january': '01', 'febrero': '02', 'march': '03', 'april': '04',
        'may': '05', 'june': '06', 'july': '07', 'august': '08',
        'september': '09', 'october': '10', 'november': '11', 'december': '12',
        'enero': '01', 'february': '02', 'marzo': '03', 'abril': '04',
        'mayo': '05', 'junio': '06', 'julio': '07', 'agosto': '08',
        'septiembre': '09', 'octubre': '10', 'noviembre': '11', 'diciembre': '12',
        'jan': '01', 'feb': '02', 'mar': '03', 'apr': '04',
        'jun': '06', 'jul': '07', 'aug': '08', 'sep': '09', 
        'oct': '10', 'nov': '11', 'dec': '12'
    }
    
    try:
        # Caso 1: Solo año (ej: "2020", "©2020", "Copyright 2020")
        year_only = re.search(r'\b(19|20)\d{2}\b', date_string)
        if year_only and len(date_string) <= 15:
            year = year_only.group()
            # Verificar si hay más información de mes
            date_lower = date_string.lower()
            for month_name, month_num in months.items():
                if month_name in date_lower:
                    return f"01/{month_num}/{year}"
            return f"01/01/{year}"
        
        # Caso 2: Formato ISO (YYYY-MM-DD)
        iso_match = re.match(r'(\d{4})-(\d{2})-(\d{2})', date_string)
        if iso_match:
            year, month, day = iso_match.groups()
            return f"{day}/{month}/{year}"
        
        # Caso 3: Formato DD/MM/YYYY o DD-MM-YYYY
        dmy_match = re.match(r'(\d{1,2})[/-](\d{1,2})[/-](\d{4})', date_string)
        if dmy_match:
            day, month, year = dmy_match.groups()
            return f"{day.zfill(2)}/{month.zfill(2)}/{year}"
        
        # Caso 4: Formato MM/DD/YYYY (americano) - asumimos si mes > 12
        mdy_match = re.match(r'(\d{1,2})[/-](\d{1,2})[/-](\d{4})', date_string)
        if mdy_match:
            first, second, year = mdy_match.groups()
            if int(first) > 12:  # Probablemente es día
                return f"{first.zfill(2)}/{second.zfill(2)}/{year}"
            else:  # Asumimos DD/MM/YYYY
                return f"{first.zfill(2)}/{second.zfill(2)}/{year}"
        
        # Caso 5: Formato con nombre de mes (ej: "March 15, 2020" o "15 de marzo de 2020")
        date_lower = date_string.lower()
        for month_name, month_num in months.items():
            if month_name in date_lower:
                year_match = re.search(r'(19|20)\d{2}', date_string)
                day_match = re.search(r'\b(\d{1,2})\b', date_string)
                if year_match:
                    year = year_match.group()
                    day = day_match.group().zfill(2) if day_match and int(day_match.group()) <= 31 else "01"
                    return f"{day}/{month_num}/{year}"
        
        # Caso 6: Solo año encontrado en string más largo
        if year_only:
            return f"01/01/{year_only.group()}"
        
        return "NO_IDENTIFICADO"
        
    except Exception as e:
        print(f"Error formateando fecha '{date_string}': {str(e)}")
        return "NO_IDENTIFICADO"


# ============================================================
# FUNCIÓN DE COMPARACIÓN Y ACTUALIZACIÓN DE METADATA
# ============================================================

async def compare_and_update_metadata(
    extracted_metadata: Dict[str, Any],
    library_id: str,
    current_data: Dict[str, Any],
    supabase: Client
) -> Dict[str, Any]:
    """
    Compara la metadata extraída con la existente y actualiza si corresponde.
    
    Reglas:
    - Si campo vacío y hay dato extraído → actualizar con dato extraído
    - Si campo vacío y NO hay dato extraído → poner "No identificado"
    - Si campo tiene dato y hay dato extraído diferente → actualizar con dato extraído
    - Si campo tiene dato y NO hay dato extraído → no tocar
    """
    try:
        updates_library = {}
        updates_documents = {}
        changes_made = []
        
        # --- Procesar TÍTULO ---
        extracted_title = extracted_metadata.get("title", "NO_ENCONTRADO")
        current_title = current_data.get("title") or ""
        
        if extracted_title != "NO_ENCONTRADO":
            # Hay título extraído
            if not current_title.strip():
                # Campo vacío → actualizar
                updates_library["title"] = extracted_title
                updates_documents["title"] = extracted_title
                changes_made.append(f"title: vacío → '{extracted_title}'")
            elif current_title.strip().lower() != extracted_title.strip().lower():
                # Campo diferente → actualizar
                updates_library["title"] = extracted_title
                updates_documents["title"] = extracted_title
                changes_made.append(f"title: '{current_title}' → '{extracted_title}'")
        else:
            # No se encontró título
            if not current_title.strip():
                updates_library["title"] = "No identificado"
                updates_documents["title"] = "No identificado"
                changes_made.append("title: vacío → 'No identificado'")
        
        # --- Procesar AUTORES ---
        extracted_authors = extracted_metadata.get("authors", ["NO_ENCONTRADO"])
        current_authors = current_data.get("authors") or []
        
        # Normalizar: si es string, convertir a lista
        if isinstance(current_authors, str):
            current_authors = [current_authors] if current_authors.strip() else []
        
        # Verificar si hay autores válidos extraídos
        valid_extracted_authors = [a for a in extracted_authors if a != "NO_ENCONTRADO"]
        
        if valid_extracted_authors:
            # Hay autores extraídos
            if not current_authors:
                # Campo vacío → actualizar
                updates_library["authors"] = valid_extracted_authors
                updates_documents["authors"] = valid_extracted_authors
                changes_made.append(f"authors: vacío → {valid_extracted_authors}")
            elif set(a.lower().strip() for a in current_authors) != set(a.lower().strip() for a in valid_extracted_authors):
                # Autores diferentes → actualizar
                updates_library["authors"] = valid_extracted_authors
                updates_documents["authors"] = valid_extracted_authors
                changes_made.append(f"authors: {current_authors} → {valid_extracted_authors}")
        else:
            # No se encontraron autores
            if not current_authors:
                updates_library["authors"] = ["No identificado"]
                updates_documents["authors"] = ["No identificado"]
                changes_made.append("authors: vacío → ['No identificado']")
        
        # --- Procesar FECHA DE PUBLICACIÓN ---
        extracted_date = extracted_metadata.get("publication_date", "NO_ENCONTRADO")
        formatted_date = format_publication_date(extracted_date)
        current_date = current_data.get("publication_date") or ""
        
        if formatted_date != "NO_IDENTIFICADO":
            # Hay fecha extraída válida
            if not current_date.strip():
                # Campo vacío → actualizar
                updates_library["publication_date"] = formatted_date
                updates_documents["publication_date"] = formatted_date
                changes_made.append(f"publication_date: vacío → '{formatted_date}'")
            elif current_date.strip() != formatted_date:
                # Fecha diferente → actualizar
                updates_library["publication_date"] = formatted_date
                updates_documents["publication_date"] = formatted_date
                changes_made.append(f"publication_date: '{current_date}' → '{formatted_date}'")
        else:
            # No se encontró fecha
            if not current_date.strip():
                updates_library["publication_date"] = "No identificado"
                updates_documents["publication_date"] = "No identificado"
                changes_made.append("publication_date: vacío → 'No identificado'")
        
        # --- Aplicar actualizaciones ---
        
        # Actualizar AIvantage_library
        if updates_library:
            supabase.table("AIvantage_library").update(
                updates_library
            ).eq("id", library_id).execute()
            print(f"AIvantage_library actualizado: {updates_library}")
        
        # Actualizar metadata en documents_v2 (todos los chunks del libro)
        docs_result = None
        if updates_documents:
            # Obtener todos los documentos con este library_id en metadata
            docs_result = supabase.table("documents_v2").select("id, metadata").filter(
                "metadata->>library_id", "eq", library_id
            ).execute()
            
            if docs_result.data:
                for doc in docs_result.data:
                    current_metadata = doc.get("metadata", {})
                    # Actualizar solo los campos que cambiaron
                    updated_metadata = {**current_metadata, **updates_documents}
                    
                    supabase.table("documents_v2").update({
                        "metadata": updated_metadata
                    }).eq("id", doc["id"]).execute()
                
                print(f"documents_v2 actualizado: {len(docs_result.data)} chunks")
        
        return {
            "success": True,
            "changes_made": changes_made,
            "updates_library": updates_library,
            "updates_documents": updates_documents,
            "total_chunks_updated": len(docs_result.data) if updates_documents and docs_result and docs_result.data else 0
        }
        
    except Exception as e:
        print(f"Error en compare_and_update_metadata: {str(e)}")
        return {
            "success": False,
            "error": str(e),
            "changes_made": []
        }


# ============================================================
# FUNCIÓN DE EXTRACCIÓN DE METADATA CON LLM
# ============================================================

async def extract_metadata_with_llm(
    chunks: List[Dict[str, Any]], 
    openai_client: OpenAI,
    file_name: str
) -> Dict[str, Any]:
    """
    Extrae metadata (título, autores, fecha de publicación) usando GPT-4o-mini.
    Analiza los primeros y últimos chunks donde típicamente está esta información.
    """
    try:
        # Obtener los primeros chunks (portada, índice, introducción)
        first_chunks = chunks[:3] if len(chunks) >= 3 else chunks
        # Obtener los últimos chunks (bibliografía, contraportada)
        last_chunks = chunks[-2:] if len(chunks) >= 2 else []
        
        # Combinar el texto relevante
        relevant_text = ""
        for chunk in first_chunks:
            relevant_text += chunk['content'][:2000] + "\n\n"
        for chunk in last_chunks:
            relevant_text += chunk['content'][:1000] + "\n\n"
        
        # Limitar el texto total
        relevant_text = relevant_text[:8000]
        
        # Prompt para extracción de metadata
        prompt = f"""Analiza el siguiente texto extraído de un documento académico y extrae la siguiente información:

1. **Título del libro/documento**: El título principal de la obra.
2. **Autores**: Lista de autores en formato "Nombre Apellido". Si hay múltiples autores, sepáralos.
3. **Fecha de publicación**: Año o fecha completa de publicación.

Nombre del archivo: {file_name}

Texto del documento:
---
{relevant_text}
---

Responde SOLO en formato JSON con esta estructura exacta:
{{
    "title": "Título del libro" o "NO_ENCONTRADO",
    "authors": ["Autor 1", "Autor 2"] o ["NO_ENCONTRADO"],
    "publication_date": "YYYY" o "DD/MM/YYYY" o "NO_ENCONTRADO"
}}

IMPORTANTE:
- Si no encuentras algún dato con certeza, usa "NO_ENCONTRADO"
- Para autores, usa el formato "Nombre Apellido" (ej: "Stephen P. Robbins")
- Para la fecha, preferiblemente extrae solo el año
- Basa tu respuesta en el contenido del texto, no en el nombre del archivo"""

        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "Eres un experto en extracción de metadata de documentos académicos. Respondes solo en JSON válido."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.1,
            max_tokens=500
        )
        
        # Parsear respuesta
        response_text = response.choices[0].message.content.strip()
        
        # Limpiar el JSON si viene con markdown
        if response_text.startswith("```"):
            response_text = response_text.split("```")[1]
            if response_text.startswith("json"):
                response_text = response_text[4:]
        response_text = response_text.strip()
        
        metadata = json.loads(response_text)
        
        print(f"Metadata extraída por LLM: {metadata}")
        return metadata
        
    except json.JSONDecodeError as e:
        print(f"Error parseando JSON de LLM: {str(e)}")
        return {
            "title": "NO_ENCONTRADO",
            "authors": ["NO_ENCONTRADO"],
            "publication_date": "NO_ENCONTRADO"
        }
    except Exception as e:
        print(f"Error en extract_metadata_with_llm: {str(e)}")
        return {
            "title": "NO_ENCONTRADO",
            "authors": ["NO_ENCONTRADO"],
            "publication_date": "NO_ENCONTRADO",
            "error": str(e)
        }


# ============================================================
# FUNCIONES DE EXTRACCIÓN
# ============================================================

def create_smart_chunks(text: str, chunk_size: int = 4000, overlap: int = 400) -> List[Dict[str, Any]]:
    """Crea chunks inteligentes del texto"""
    if not text or len(text) <= chunk_size:
        return [{
            'content': text,
            'chunk_index': 0,
            'total_chunks': 1,
            'start_position': 0,
            'end_position': len(text) if text else 0
        }]
    
    chunks = []
    paragraphs = text.split('\n\n')
    current_chunk = ""
    chunk_index = 0
    start_position = 0
    
    for paragraph in paragraphs:
        if len(current_chunk) + len(paragraph) + 2 > chunk_size:
            if current_chunk:
                chunks.append({
                    'content': current_chunk.strip(),
                    'chunk_index': chunk_index,
                    'start_position': start_position,
                    'end_position': start_position + len(current_chunk)
                })
                chunk_index += 1
                
                lines = current_chunk.split('\n')
                overlap_lines = lines[-5:] if len(lines) > 5 else lines
                overlap_text = '\n'.join(overlap_lines)
                
                start_position = start_position + len(current_chunk) - len(overlap_text)
                current_chunk = overlap_text + "\n\n" + paragraph
            else:
                if len(paragraph) > chunk_size:
                    sentences = paragraph.split('. ')
                    for sentence in sentences:
                        if len(current_chunk) + len(sentence) + 2 > chunk_size:
                            chunks.append({
                                'content': current_chunk.strip(),
                                'chunk_index': chunk_index,
                                'start_position': start_position,
                                'end_position': start_position + len(current_chunk)
                            })
                            chunk_index += 1
                            start_position += len(current_chunk)
                            current_chunk = sentence + ". "
                        else:
                            current_chunk += sentence + ". "
                else:
                    current_chunk = paragraph
        else:
            if current_chunk:
                current_chunk += "\n\n" + paragraph
            else:
                current_chunk = paragraph
    
    if current_chunk:
        chunks.append({
            'content': current_chunk.strip(),
            'chunk_index': chunk_index,
            'start_position': start_position,
            'end_position': len(text)
        })
    
    for chunk in chunks:
        chunk['total_chunks'] = len(chunks)
    
    return chunks


async def extract_pdf(content: bytes) -> Dict[str, Any]:
    """Extrae texto de PDF (versión sin progreso, para compatibilidad)"""
    result = {'text': '', 'pages': 0}
    
    try:
        pdf_file = io.BytesIO(content)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        result['pages'] = len(pdf_reader.pages)
        
        text_parts = []
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"--- Página {page_num + 1} ---\n{page_text}")
        
        result['text'] = '\n\n'.join(text_parts)
        
        # Si no hay texto, intentar con pdfplumber
        if not result['text'].strip() or len(result['text']) < 100:
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            with pdfplumber.open(tmp_path) as pdf:
                result['pages'] = len(pdf.pages)
                text_parts = []
                for i, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(f"--- Página {i + 1} ---\n{page_text}")
                result['text'] = '\n\n'.join(text_parts)
            
            os.unlink(tmp_path)
        
        # Si aún no hay texto, intentar OCR
        if not result['text'].strip() or len(result['text']) < 100:
            result['text'] = await extract_pdf_with_ocr(content)
            
    except Exception as e:
        print(f"Error extrayendo PDF: {e}")
        result['text'] = await extract_pdf_with_ocr(content)
    
    return result


async def extract_pdf_with_ocr(content: bytes) -> str:
    """Extrae texto de PDF usando OCR (versión sin progreso)"""
    try:
        images = convert_from_bytes(content, dpi=200)
        text_parts = []
        
        for i, image in enumerate(images):
            text = pytesseract.image_to_string(image, lang='spa+eng')
            text_parts.append(f"--- Página {i + 1} (OCR) ---\n{text}")
        
        return '\n\n'.join(text_parts)
    except Exception as e:
        return f"Error en OCR: {str(e)}"


async def extract_docx(content: bytes) -> Dict[str, Any]:
    """Extrae texto de archivos Word"""
    doc_stream = io.BytesIO(content)
    doc = Document(doc_stream)
    
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    
    return {
        'text': '\n\n'.join(text_parts),
        'paragraphs': len(doc.paragraphs)
    }


async def extract_excel(content: bytes) -> Dict[str, Any]:
    """Extrae datos de archivos Excel"""
    excel_stream = io.BytesIO(content)
    workbook = openpyxl.load_workbook(excel_stream, data_only=True)
    
    all_text = []
    for sheet_name in workbook.sheetnames:
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows():
            row_values = [str(cell.value) for cell in row if cell.value is not None]
            if row_values:
                all_text.append(' | '.join(row_values))
    
    return {
        'text': '\n'.join(all_text),
        'total_sheets': len(workbook.sheetnames)
    }


async def extract_pptx(content: bytes) -> Dict[str, Any]:
    """Extrae texto de PowerPoint"""
    pptx_stream = io.BytesIO(content)
    presentation = Presentation(pptx_stream)
    
    all_text = []
    for i, slide in enumerate(presentation.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text)
        
        if slide_text:
            all_text.append(f"--- Diapositiva {i + 1} ---\n" + '\n'.join(slide_text))
    
    return {
        'text': '\n\n'.join(all_text),
        'total_slides': len(presentation.slides)
    }


async def extract_image(content: bytes) -> Dict[str, Any]:
    """Extrae texto de imágenes usando OCR"""
    image = Image.open(io.BytesIO(content))
    text = pytesseract.image_to_string(image, lang='spa+eng')
    return {'text': text.strip()}


async def extract_text(content: bytes) -> Dict[str, Any]:
    """Extrae texto de archivos de texto plano"""
    try:
        text = content.decode('utf-8')
    except UnicodeDecodeError:
        text = content.decode('latin-1')
    return {'text': text}


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
